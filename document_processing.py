import PyPDF2
from pptx import Presentation

def process_document(file_path):
    ext = file_path.split('.')[-1].lower()
    if ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['ppt', 'pptx']:
        return extract_text_from_ppt(file_path)
    else:
        raise ValueError("Unsupported file type")

def extract_text_from_pdf(file_path):
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text
